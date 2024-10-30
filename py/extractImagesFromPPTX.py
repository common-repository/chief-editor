#!/usr/bin/env python
# -*- coding: utf-8 -*-
import sys
import os
from pptx import Presentation;
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import StringIO

def main():
    input = sys.argv[1]
    output = sys.argv[2]
    print("Extract")
    print(input)
    print("in")
    
    output = os.path.normpath(output)
    print(output)
    # extract text and write images in /tmp/img_dir
    #text = docxpy.process(input, output)  
    idx = 1
    prs = Presentation(input);
    for slide in prs.slides:
        #print("slide")
        for shape in slide.shapes:
            #print("shape")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                #print(shape._pic.nvPicPr.cNvPr.get('descr'))
                '''
                print(shape.shape_type)
                print(shape.image.filename)
                print(shape.image.ext)   
                print(shape.image)'''
                imageBytes = shape.image.blob
                #print(imageBytes)
                im = Image.open(StringIO.StringIO(imageBytes))
                outFilename = os.path.normpath(output + '/' +'__'+str(idx)+'__'+ shape.image.filename)# + shape.image.ext
                print(outFilename)
                im.save(outFilename)
                idx = idx + 1
            #    print(shape.text)  

if __name__== "__main__":
  main()